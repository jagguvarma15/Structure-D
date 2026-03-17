"""Office document parsing (DOCX, XLSX, PPTX)."""

from __future__ import annotations

import asyncio
from pathlib import Path
from typing import Any

import structlog

from structure_d.ingestion.base import BaseParser
from structure_d.schemas.base import DocumentFormat, DocumentMetadata, ParsedDocument

logger = structlog.get_logger(__name__)


class DocxParser(BaseParser):
    """Parse .docx files using python-docx.

    Extracts paragraphs (with heading markers), table content, core document
    properties (author, title, dates) and an estimated page count derived from
    explicit page-break elements in the XML.
    """

    supported_extensions = [".docx"]

    async def parse(self, file_path: Path, **kwargs: object) -> ParsedDocument:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._parse_sync, file_path)

    def _parse_sync(self, file_path: Path) -> ParsedDocument:
        try:
            from docx import Document
            from docx.oxml.ns import qn
            from docx.table import Table as DocxTable
            from docx.text.paragraph import Paragraph as DocxParagraph
        except ImportError as exc:
            raise ImportError(
                "python-docx is required for DOCX parsing. "
                "Install with: pip install 'structure-d[ingestion]'"
            ) from exc

        doc = Document(str(file_path))

        # ── Iterate body elements in document order ───────────────────────────
        text_parts: list[str] = []
        tables: list[dict[str, Any]] = []
        table_index = 0

        for item in self._iter_block_items(doc):
            if isinstance(item, DocxParagraph):
                line = self._format_paragraph(item)
                if line:
                    text_parts.append(line)
            else:  # DocxTable
                rows = self._extract_table_rows(item)
                tables.append({"index": table_index, "data": rows})
                table_index += 1
                # Weave table text into the main content flow
                for row in rows:
                    row_text = "\t".join(row)
                    if row_text.strip():
                        text_parts.append(row_text)

        full_text = "\n\n".join(text_parts)

        # ── Estimate page count from explicit page-break elements ─────────────
        page_breaks = [
            br
            for br in doc.element.body.iter(qn("w:br"))
            if br.get(qn("w:type")) == "page"
        ]
        page_count = len(page_breaks) + 1

        # ── Core document properties ──────────────────────────────────────────
        cp = doc.core_properties
        extra: dict[str, Any] = {}
        if cp.title:
            extra["title"] = cp.title
        if cp.author:
            extra["author"] = cp.author
        if cp.created:
            extra["created"] = cp.created.isoformat()
        if cp.modified:
            extra["modified"] = cp.modified.isoformat()
        if cp.last_modified_by:
            extra["last_modified_by"] = cp.last_modified_by

        metadata = DocumentMetadata(
            filename=file_path.name,
            source=str(file_path),
            file_extension=file_path.suffix.lower(),
            format=DocumentFormat.DOCX,
            file_size_bytes=file_path.stat().st_size,
            page_count=page_count,
            extra=extra,
        )

        logger.debug(
            "docx.parsed",
            filename=file_path.name,
            page_count=page_count,
            tables=len(tables),
            chars=len(full_text),
        )

        return ParsedDocument(
            metadata=metadata,
            text=full_text,
            pages=[full_text],
            tables=tables,
        )

    @staticmethod
    def _iter_block_items(doc: object):
        """Yield Paragraph and Table objects in document body order."""
        from docx.oxml.ns import qn
        from docx.table import Table as DocxTable
        from docx.text.paragraph import Paragraph as DocxParagraph

        for child in doc.element.body.iterchildren():
            if child.tag == qn("w:p"):
                yield DocxParagraph(child, doc)
            elif child.tag == qn("w:tbl"):
                yield DocxTable(child, doc)

    @staticmethod
    def _format_paragraph(para: object) -> str:
        """Return paragraph text prefixed with Markdown heading markers."""
        text = para.text.strip()
        if not text:
            return ""
        style_name: str = para.style.name if para.style else ""
        if style_name.startswith("Heading 1"):
            return f"# {text}"
        if style_name.startswith("Heading 2"):
            return f"## {text}"
        if style_name.startswith("Heading 3"):
            return f"### {text}"
        return text

    @staticmethod
    def _extract_table_rows(table: object) -> list[list[str]]:
        """Return all cell text from *table* as a 2-D list of strings."""
        return [[cell.text.strip() for cell in row.cells] for row in table.rows]


class XlsxParser(BaseParser):
    """Parse .xlsx files using openpyxl."""

    supported_extensions = [".xlsx"]

    async def parse(self, file_path: Path, **kwargs: object) -> ParsedDocument:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._parse_sync, file_path)

    def _parse_sync(self, file_path: Path) -> ParsedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        pages: list[str] = []
        tables: list[dict[str, Any]] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c is not None else "" for c in row])
            tables.append({"sheet": sheet_name, "data": rows})
            page_text = "\n".join(["\t".join(r) for r in rows])
            pages.append(page_text)

        wb.close()
        text = "\n\n".join(pages)

        metadata = DocumentMetadata(
            filename=file_path.name,
            source=str(file_path),
            file_extension=".xlsx",
            file_size_bytes=file_path.stat().st_size,
            page_count=len(pages),
        )

        return ParsedDocument(metadata=metadata, text=text, pages=pages, tables=tables)


class PptxParser(BaseParser):
    """Parse .pptx files using python-pptx."""

    supported_extensions = [".pptx"]

    async def parse(self, file_path: Path, **kwargs: object) -> ParsedDocument:
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._parse_sync, file_path)

    def _parse_sync(self, file_path: Path) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        pages: list[str] = []

        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
            pages.append("\n".join(parts))

        text = "\n\n".join(pages)

        metadata = DocumentMetadata(
            filename=file_path.name,
            source=str(file_path),
            file_extension=".pptx",
            file_size_bytes=file_path.stat().st_size,
            page_count=len(pages),
        )

        return ParsedDocument(metadata=metadata, text=text, pages=pages)
